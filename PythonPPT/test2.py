from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Introduction
slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Visualization Types in Banking"
subtitle = slide.placeholders[1]
subtitle.text = "Examples and Use Cases"

# Slide 2: Bar Chart
slide_layout = prs.slide_layouts[5]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Bar Chart"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Useful for comparing discrete categories. Examples:\n- Age distribution by branch\n- Loan types by branch"

# Slide 3: Line Graph
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Line Graph"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Ideal for showing trends over time. Examples:\n- Average account balance over months\n- Customer growth over years"

# Slide 4: Scatter Plot
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Scatter Plot"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Used to visualize relationships between two variables. Examples:\n- Account balance vs. number of transactions\n- Age vs. income"

# Slide 5: Pie Chart
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Pie Chart"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Illustrates proportions of a whole. Examples:\n- Loan types distribution within a branch\n- Expenses breakdown by category"

# Slide 6: Histogram
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Histogram"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Displays distribution of a continuous variable. Examples:\n- Income distribution within an age group\n- Transaction amounts distribution"

# Slide 7: Heatmap
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Heatmap"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Effective for visualizing relationship between two categorical variables. Examples:\n- ATM transactions frequency by time of day and day of week\n- Loan approval rates by customer segment"

# Slide 8: Box Plot
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Box Plot"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Displays distribution and outliers of a continuous variable. Examples:\n- Account balance distribution by account type\n- Loan amount distribution by loan type"

# Save the presentation
prs.save('banking_visualization_examples.pptx')
