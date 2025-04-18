from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import numpy as np

# Create a presentation object
prs = Presentation()

# Slide 1: Introduction
slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Visualization Types in Banking"
subtitle = slide.placeholders[1]
subtitle.text = "Examples and Use Cases"

# Slide 2: Bar Chart
slide_layout = prs.slide_layouts[5]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Bar Chart"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Useful for comparing discrete categories. Examples:\n- Age distribution by branch\n- Loan types by branch"

# Generate sample data for the bar chart
age_groups = ['18-25', '26-35', '36-45']
branches = ['Branch A', 'Branch B', 'Branch C']
data = np.random.randint(10, 100, size=(len(branches), len(age_groups)))

# Create a bar chart
fig, ax = plt.subplots()
bar_width = 0.2
for i, branch in enumerate(branches):
    ax.bar(np.arange(len(age_groups)) + i * bar_width, data[i], bar_width, label=branch)
ax.set_xticks(np.arange(len(age_groups)) + bar_width / 2)
ax.set_xticklabels(age_groups)
ax.legend()
ax.set_xlabel('Age Groups')
ax.set_ylabel('Number of Customers')
ax.set_title('Example Bar Chart')
plt.close(fig)

# Save the bar chart as an image
fig.savefig('bar_chart.png')

# Add the bar chart image to the slide
left = Inches(1)
top = Inches(3)
slide.shapes.add_picture('bar_chart.png', left, top, width=Inches(8))

# Slide 3: Line Graph
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Line Graph"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Ideal for showing trends over time. Examples:\n- Average account balance over months\n- Customer growth over years"

# Generate sample data for the line graph
months = ['Jan', 'Feb', 'Mar', 'Apr', 'May']
branches = ['Branch A', 'Branch B', 'Branch C']
data = np.random.randint(1000, 5000, size=(len(branches), len(months)))

# Create a line graph
fig, ax = plt.subplots()
for i, branch in enumerate(branches):
    ax.plot(months, data[i], label=branch)
ax.legend()
ax.set_xlabel('Months')
ax.set_ylabel('Average Account Balance')
ax.set_title('Example Line Graph')
plt.close(fig)

# Save the line graph as an image
fig.savefig('line_graph.png')

# Add the line graph image to the slide
left = Inches(1)
top = Inches(3)
slide.shapes.add_picture('line_graph.png', left, top, width=Inches(8))

# Slide 4: Scatter Plot
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Scatter Plot"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Used to visualize relationships between two variables. Examples:\n- Account balance vs. number of transactions\n- Age vs. income"

# Generate sample data for the scatter plot
account_balance = np.random.uniform(1000, 10000, size=100)
num_transactions = np.random.randint(1, 100, size=100)

# Create a scatter plot
fig, ax = plt.subplots()
ax.scatter(account_balance, num_transactions, alpha=0.5)
ax.set_xlabel('Account Balance')
ax.set_ylabel('Number of Transactions')
ax.set_title('Example Scatter Plot')
plt.close(fig)

# Save the scatter plot as an image
fig.savefig('scatter_plot.png')

# Add the scatter plot image to the slide
left = Inches(1)
top = Inches(3)
slide.shapes.add_picture('scatter_plot.png', left, top, width=Inches(8))

# Slide 5: Pie Chart
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Pie Chart"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Illustrates proportions of a whole. Examples:\n- Loan types distribution within a branch\n- Expenses breakdown by category"

# Generate sample data for the pie chart
loan_types = ['Mortgage', 'Personal Loan', 'Auto Loan']
loan_counts = np.random.randint(100, 1000, size=len(loan_types))

# Create a pie chart
fig, ax = plt.subplots()
ax.pie(loan_counts, labels=loan_types, autopct='%1.1f%%', startangle=140)
ax.axis('equal')  # Equal aspect ratio ensures that pie is drawn as a circle.
ax.set_title('Example Pie Chart')
plt.close(fig)

# Save the pie chart as an image
fig.savefig('pie_chart.png')

# Add the pie chart image to the slide
left = Inches(1)
top = Inches(3)
slide.shapes.add_picture('pie_chart.png', left, top, width=Inches(8))

# Slide 6: Histogram
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Histogram"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Displays distribution of a continuous variable. Examples:\n- Income distribution within an age group\n- Transaction amounts distribution"

# Generate sample data for the histogram
income = np.random.normal(50000, 20000, 1000)

# Create a histogram
fig, ax = plt.subplots()
ax.hist(income, bins=30, edgecolor='black', alpha=0.7)
ax.set_xlabel('Income')
ax.set_ylabel('Frequency')
ax.set_title('Example Histogram')
plt.close(fig)

# Save the histogram as an image
fig.savefig('histogram.png')

# Add the histogram image to the slide
left = Inches(1)
top = Inches(3)
slide.shapes.add_picture('histogram.png', left, top, width=Inches(8))

# Slide 7: Heatmap
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Heatmap"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Effective for visualizing relationship between two categorical variables. Examples:\n- ATM transactions frequency by time of day and day of week\n- Loan approval rates by customer segment"

# Generate sample data for the heatmap
days = ['Mon', 'Tue', 'Wed', 'Thu', 'Fri', 'Sat', 'Sun']
hours = ['12am', '1am', '2am', '3am', '4am', '5am', '6am', '7am', '8am', '9am', '10am', '11am', '12pm', '1pm', '2pm', '3pm', '4pm', '5pm', '6pm', '7pm', '8pm', '9pm', '10pm', '11pm']
heatmap_data = np.random.randint(0, 100, size=(len(days), len(hours)))

# Create a heatmap
fig, ax = plt.subplots()
cax = ax.imshow(heatmap_data, cmap='hot', interpolation='nearest')
ax.set_xticks(np.arange(len(hours)))
ax.set_yticks(np.arange(len(days)))
ax.set_xticklabels(hours)
ax.set_yticklabels(days)
plt.colorbar(cax)
ax.set_title('Example Heatmap')
plt.close(fig)

# Save the heatmap as an image
fig.savefig('heatmap.png')

# Add the heatmap image to the slide
left = Inches(1)
top = Inches(3)
slide.shapes.add_picture('heatmap.png', left, top, width=Inches(8))

# Slide 8: Box Plot
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Box Plot"
content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
content.text = "Displays distribution and outliers of a continuous variable. Examples:\n- Account balance distribution by account type\n- Loan amount distribution by loan type"

# Generate sample data for the box plot
account_types = ['Savings', 'Checking', 'Investment']
account_balances = [np.random.normal(5000, 1000, 100), np.random.normal(3000, 500, 100), np.random.normal(10000, 2000, 100)]

# Create a box plot
fig, ax = plt.subplots()
ax.boxplot(account_balances, labels=account_types)
ax.set_xlabel('Account Type')
ax.set_ylabel('Account Balance')
ax.set_title('Example Box Plot')
plt.close(fig)

# Save the box plot as an image
fig.savefig('box_plot.png')

# Add the box plot image to the slide
left = Inches(1)
top = Inches(3)
slide.shapes.add_picture('box_plot.png', left, top, width=Inches(8))

# Save the presentation
prs.save('banking_visualization_examplesFinal.pptx')
